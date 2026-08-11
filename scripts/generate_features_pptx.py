"""Generate a simple HRMS Pro features presentation (.pptx)."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

FOREST = RGBColor(0x1B, 0x4D, 0x3E)
TEAL = RGBColor(0x2A, 0x6F, 0x5F)
ACCENT = RGBColor(0xC4, 0xA3, 0x5A)
LIGHT = RGBColor(0xF5, 0xF7, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1A, 0x1F, 0x1E)
MUTED = RGBColor(0x5A, 0x66, 0x62)


def set_run(run, size=18, bold=False, color=DARK, font='Calibri'):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_bg(slide, color=LIGHT):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    sp_tree = slide.shapes._spTree
    sp = shape._element
    sp_tree.remove(sp)
    sp_tree.insert(2, sp)


def add_bar(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(
    slide,
    left,
    top,
    width,
    height,
    text,
    size=18,
    bold=False,
    color=DARK,
    align=PP_ALIGN.LEFT,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color)
    return box


def add_bullets(slide, left, top, width, height, items, size=18, color=DARK, spacing=10):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(spacing)
        run = p.add_run()
        run.text = '•  ' + item
        set_run(run, size=size, color=color)
    return box


def add_card(slide, left, top, width, height, title, lines):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(0xD5, 0xE0, 0xDB)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.12), height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()
    add_text(
        slide,
        left + Inches(0.3),
        top + Inches(0.18),
        width - Inches(0.4),
        Inches(0.4),
        title,
        size=16,
        bold=True,
        color=FOREST,
    )
    add_bullets(
        slide,
        left + Inches(0.25),
        top + Inches(0.55),
        width - Inches(0.4),
        height - Inches(0.7),
        lines,
        size=13,
        color=MUTED,
        spacing=6,
    )


def section_header(slide, title, subtitle=''):
    add_bg(slide, LIGHT)
    add_bar(slide, 0, 0, prs.slide_width, Inches(0.08), FOREST)
    add_bar(slide, 0, Inches(7.42), prs.slide_width, Inches(0.08), FOREST)
    add_text(
        slide,
        Inches(0.7),
        Inches(0.35),
        Inches(11),
        Inches(0.55),
        title,
        size=32,
        bold=True,
        color=FOREST,
    )
    if subtitle:
        add_text(
            slide,
            Inches(0.7),
            Inches(0.9),
            Inches(11.5),
            Inches(0.4),
            subtitle,
            size=16,
            color=MUTED,
        )
    add_bar(slide, Inches(0.7), Inches(1.35), Inches(1.2), Inches(0.06), ACCENT)


# 1. Title
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, FOREST)
add_bar(slide, 0, Inches(5.9), prs.slide_width, Inches(1.6), TEAL)
add_text(
    slide, Inches(0.9), Inches(2.0), Inches(11), Inches(0.5),
    'HRMS PRO', size=20, bold=True, color=ACCENT,
)
add_text(
    slide, Inches(0.9), Inches(2.5), Inches(11.5), Inches(1.2),
    'Human Resource & Payroll System', size=40, bold=True, color=WHITE,
)
add_text(
    slide, Inches(0.9), Inches(3.8), Inches(11), Inches(0.6),
    'A simple walkthrough of every feature in the platform',
    size=20, color=RGBColor(0xC8, 0xD9, 0xD3),
)
add_text(
    slide, Inches(0.9), Inches(6.25), Inches(11), Inches(0.4),
    'Multi-tenant  ·  Secure  ·  AI-assisted  ·  End-to-end HR',
    size=16, color=WHITE,
)
add_text(
    slide, Inches(0.9), Inches(6.7), Inches(11), Inches(0.35),
    'Presentation Overview  |  2026', size=14, color=RGBColor(0xB0, 0xC8, 0xC0),
)

# 2. Agenda
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 'Agenda', 'What we will cover today')
add_bullets(
    slide, Inches(0.9), Inches(1.7), Inches(5.8), Inches(5),
    [
        '1.  What is HRMS Pro?',
        '2.  Who uses it & how tenants work',
        '3.  People modules (Employees, Recruitment, Casuals)',
        '4.  Daily operations (Leave, Attendance, Payroll)',
    ],
    size=20, spacing=18,
)
add_bullets(
    slide, Inches(7.0), Inches(1.7), Inches(5.5), Inches(5),
    [
        '5.  Growth & governance (Performance, Relations, Surveys)',
        '6.  AI, Reports & Integrations',
        '7.  Security, roles & settings',
        '8.  Key benefits & closing',
    ],
    size=20, spacing=18,
)

# 3. What is it
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 'What is HRMS Pro?', 'One platform for the full employee lifecycle')
add_bullets(
    slide, Inches(0.9), Inches(1.7), Inches(11.5), Inches(5),
    [
        'A complete Human Resource Management System for companies of any size',
        'Covers hiring → onboarding → attendance → leave → payroll → performance',
        'Each company (tenant) has its own secure workspace and data',
        'Web dashboard for HR, managers, payroll officers, and employees',
        'REST API ready for mobile apps and external systems',
        'Smart AI helpers for screening, anomalies, and leave decisions',
    ],
    size=20, spacing=14,
)

# 4. Who uses it
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 'Who Uses the System?', 'Role-based access — everyone sees what they need')
roles = [
    ('Super Admin', ['Full control', 'Enable/disable roles', 'Company settings']),
    ('HR Admin', ['Employees & hiring', 'Leave policies', 'Org structure']),
    ('Payroll Officer', ['Payroll runs', 'Payslips & loans', 'Compliance checks']),
    ('Managers', ['Approve leave', 'Team performance', 'Attendance oversight']),
    ('Employees', ['Self-service portal', 'Apply leave', 'View payslips']),
    ('Auditor', ['Read-only access', 'Audit trails', 'Compliance review']),
]
for i, (title, lines) in enumerate(roles):
    r, c = divmod(i, 3)
    left = Inches(0.7) + Inches(c * 4.1)
    top = Inches(1.65) + Inches(r * 2.55)
    add_card(slide, left, top, Inches(3.9), Inches(2.35), title, lines)

# 5. Multi-tenant
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 'Multi-Tenant by Design', 'Every company is isolated and secure')
add_bullets(
    slide, Inches(0.9), Inches(1.7), Inches(11.5), Inches(5),
    [
        'Each tenant has a unique portal URL, e.g. /t/acme-corp/',
        'Company data never mixes with another organization',
        'Branches, departments, and designations per company',
        'Super Admin can turn roles on/off to match that tenant’s needs',
        'Ideal for SaaS: onboard many companies on one platform',
    ],
    size=20, spacing=16,
)

# 6. Feature map
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 'Feature Map', 'All major modules at a glance')
modules = [
    ('People', ['Employees', 'Recruitment', 'Casuals']),
    ('Operations', ['Leave', 'Attendance', 'Payroll']),
    ('Growth', ['Performance', 'Surveys', 'Feedback']),
    ('Governance', ['Relations', 'Disciplinary', 'Roles & Settings']),
    ('Intelligence', ['AI Assistant', 'Reports', 'Integrations']),
]
for i, (title, lines) in enumerate(modules):
    left = Inches(0.55) + Inches(i * 2.5)
    add_card(slide, left, Inches(1.8), Inches(2.35), Inches(4.8), title, lines)

# 7. Employees
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, '1. Employee Management', 'Your central people directory')
add_bullets(
    slide, Inches(0.9), Inches(1.7), Inches(6), Inches(5),
    [
        'Create and manage employee profiles',
        'Track ID, contact, status, join date',
        'Link to department, designation, branch',
        'Assign reporting managers',
        'Soft-delete keeps history safe',
        'Quick search and detail panel view',
    ],
    size=19, spacing=12,
)
add_card(
    slide, Inches(7.3), Inches(1.8), Inches(5.2), Inches(4.5), 'Why it matters',
    [
        'Single source of truth for staff data',
        'Feeds leave, payroll, and performance',
        'Supports permanent & contract types',
        'HR can update records in seconds',
    ],
)

# 8. Recruitment
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, '2. Recruitment & Onboarding', 'From job posting to hired employee')
add_card(
    slide, Inches(0.7), Inches(1.7), Inches(3.9), Inches(4.8), 'Job Postings',
    ['Create open roles', 'Set department & openings', 'Track draft / open / closed', 'Closing dates & requirements'],
)
add_card(
    slide, Inches(4.8), Inches(1.7), Inches(3.9), Inches(4.8), 'Applicants',
    ['Capture candidate details', 'Link to job posting', 'AI resume score', 'Track hiring stage'],
)
add_card(
    slide, Inches(8.9), Inches(1.7), Inches(3.7), Inches(4.8), 'Interviews',
    ['Schedule interviews', 'Record ratings', 'Mark completed', 'Move to offer / hire'],
)

# 9. Leave
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, '3. Leave Management', 'Apply, approve, and track time off')
add_bullets(
    slide, Inches(0.9), Inches(1.7), Inches(6), Inches(5),
    [
        'Employees submit leave with reason & dates',
        'System checks leave balance before request',
        'Managers / HR approve or reject with notes',
        'Multi-level approval when required',
        'Balances update automatically',
        'AI can summarize pending leave for managers',
    ],
    size=19, spacing=12,
)
add_card(
    slide, Inches(7.3), Inches(1.8), Inches(5.2), Inches(4.5), 'Simple flow',
    [
        '1. Employee applies',
        '2. Balance is checked',
        '3. Supervisor reviews',
        '4. Approve or reject',
        '5. Calendar & balance update',
    ],
)

# 10. Attendance
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, '4. Time & Attendance', 'Know who is present, late, or overtime')
add_bullets(
    slide, Inches(0.9), Inches(1.7), Inches(11.5), Inches(5),
    [
        'Daily attendance records per employee',
        'Clock-in / clock-out from the portal',
        'Detect late arrivals and calculate overtime hours',
        'Compare against shifts and work schedules',
        'Supports future device sync (biometric / QR / GPS)',
        'Attendance data feeds into payroll calculations',
    ],
    size=20, spacing=14,
)

# 11. Payroll
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, '5. Payroll & Compliance', 'Run payroll with confidence')
add_card(
    slide, Inches(0.7), Inches(1.7), Inches(3.9), Inches(4.8), 'Payroll Runs',
    ['Create pay periods', 'Process employee pay', 'Gross → deductions → net', 'Review and approve'],
)
add_card(
    slide, Inches(4.8), Inches(1.7), Inches(3.9), Inches(4.8), 'Payslips',
    ['Per-employee payslips', 'Gross and net amounts', 'AI anomaly flags', 'Ready for download'],
)
add_card(
    slide, Inches(8.9), Inches(1.7), Inches(3.7), Inches(4.8), 'Loans',
    ['Staff loan records', 'Monthly installments', 'Status tracking', 'Auto payroll deduction'],
)

# 12. Performance
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, '6. Performance Management', 'Goals, reviews, and growth')
add_bullets(
    slide, Inches(0.9), Inches(1.7), Inches(11.5), Inches(5),
    [
        'Assign tasks and goals with due dates and points',
        'Track progress and completion status',
        'Support review cycles: self-assessment → manager review',
        'Capture ratings and feedback over time',
        'Helps managers coach teams with clear targets',
    ],
    size=20, spacing=16,
)

# 13. Relations & Disciplinary
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, '7. Relations & Disciplinary', 'Recognize good work and handle issues fairly')
add_card(
    slide, Inches(0.7), Inches(1.7), Inches(5.8), Inches(4.8), 'Employee Relations',
    [
        'Record awards and recognition',
        'Track employee relations cases',
        'Build a positive workplace culture',
        'Keep history linked to the employee',
    ],
)
add_card(
    slide, Inches(6.8), Inches(1.7), Inches(5.8), Inches(4.8), 'Disciplinary',
    [
        'Log incidents and investigations',
        'Warnings, hearings, outcomes',
        'Clear severity-based process',
        'Protects both company and staff',
    ],
)

# 14. Casuals & Surveys
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, '8. Casuals & Surveys', 'Flexible workers and honest feedback')
add_card(
    slide, Inches(0.7), Inches(1.7), Inches(5.8), Inches(4.8), 'Casuals Management',
    [
        'Manage temporary / casual workers',
        'Separate from permanent payroll staff',
        'Supervisor oversight for casual teams',
        'Useful for seasonal or project work',
    ],
)
add_card(
    slide, Inches(6.8), Inches(1.7), Inches(5.8), Inches(4.8), 'Feedback & Surveys',
    [
        'Create company surveys',
        'Anonymous options for honesty',
        'Set start and end dates',
        'Collect staff sentiment easily',
    ],
)

# 15. AI
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, '9. AI Features', 'Smart helpers that save HR and managers time')
add_bullets(
    slide, Inches(0.9), Inches(1.7), Inches(11.5), Inches(5),
    [
        'Resume screening & candidate scoring in recruitment',
        'Payroll anomaly detection (unusual payslips flagged)',
        'Attendance anomaly checks',
        'Leave duty assistant: summarize pending leave',
        'Draft approve / reject notes for managers',
        'AI chat insights scoped safely to your company only',
    ],
    size=20, spacing=14,
)

# 16. Reports & Integrations
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, '10. Reports & Integrations', 'See the big picture and connect systems')
add_card(
    slide, Inches(0.7), Inches(1.7), Inches(5.8), Inches(4.8), 'Reports & Insights',
    [
        'Dashboard overview of key HR metrics',
        'Module reports for leave, payroll, people',
        'Export-ready operational visibility',
        'Helps leaders make faster decisions',
    ],
)
add_card(
    slide, Inches(6.8), Inches(1.7), Inches(5.8), Inches(4.8), 'Integrations',
    [
        'REST API for external apps',
        'OAuth login (Google / Microsoft / GitHub)',
        'Background jobs via Celery',
        'Ready for biometrics & payroll banks',
    ],
)

# 17. Security
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, '11. Security & Access Control', 'Enterprise-ready protection')
add_bullets(
    slide, Inches(0.9), Inches(1.7), Inches(11.5), Inches(5),
    [
        'Secure login with email or phone',
        'Optional MFA (one-time codes)',
        'Account lockout after failed attempts',
        'Role-Based Access Control (RBAC)',
        'Tenant isolation — no cross-company data leaks',
        'Audit logs track important actions',
        'Super Admin can customize which roles a tenant uses',
    ],
    size=19, spacing=12,
)

# 18. Settings
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, '12. System Settings', 'Configure the tenant to match your organization')
add_bullets(
    slide, Inches(0.9), Inches(1.7), Inches(11.5), Inches(5),
    [
        'Manage users and assign roles',
        'Role Catalog: enable or remove roles per tenant needs',
        'Organize branches, departments, and job titles',
        'Company profile and operational preferences',
        'Keep structure clean as the business grows',
    ],
    size=20, spacing=16,
)

# 19. Benefits
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 'Why HRMS Pro?', 'Clear business value')
benefits = [
    ('One system', ['No more scattered Excel files and tools']),
    ('Faster HR', ['Hire, approve leave, and run payroll quicker']),
    ('Fewer errors', ['Balances, deductions, and AI checks']),
    ('Better control', ['Roles, tenants, and audit trails']),
    ('Scalable', ['Works for one company or many']),
    ('Future-ready', ['API + AI + integrations']),
]
for i, (t, d) in enumerate(benefits):
    r, c = divmod(i, 3)
    left = Inches(0.7) + Inches(c * 4.1)
    top = Inches(1.7) + Inches(r * 2.5)
    add_card(slide, left, top, Inches(3.9), Inches(2.25), t, d)

# 20. Closing
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, FOREST)
add_bar(slide, 0, Inches(5.9), prs.slide_width, Inches(1.6), TEAL)
add_text(
    slide, Inches(0.9), Inches(2.2), Inches(11.5), Inches(1),
    'Thank You', size=44, bold=True, color=WHITE,
)
add_text(
    slide, Inches(0.9), Inches(3.4), Inches(11.5), Inches(0.7),
    'Questions? Let’s walk through a live demo.',
    size=22, color=RGBColor(0xC8, 0xD9, 0xD3),
)
add_text(
    slide, Inches(0.9), Inches(6.3), Inches(11.5), Inches(0.5),
    'HRMS Pro  ·  Human Resource & Payroll System', size=16, color=WHITE,
)

out = r'd:\Projects Running\HR & Payroll System\docs\HRMS-Pro-Features-Presentation.pptx'
prs.save(out)
print('Saved:', out)
print('Slides:', len(prs.slides))
