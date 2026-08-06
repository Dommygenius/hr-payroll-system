"""
Generate HRMS Pro PowerPoint presentation with flowcharts and DB design.
Run: python docs/generate_presentation.py
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

OUTPUT = Path(__file__).parent / 'HRMS-System-Design.pptx'

# Brand colors
PRIMARY = RGBColor(0x3B, 0x82, 0xF6)
DARK = RGBColor(0x0F, 0x17, 0x2A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x64, 0x74, 0x8B)
LIGHT_BG = RGBColor(0xF1, 0xF5, 0xF9)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
ORANGE = RGBColor(0xEA, 0xB3, 0x08)
RED = RGBColor(0xEF, 0x44, 0x44)


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_slide(prs, title, subtitle=''):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK)
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.08))
    box.fill.solid()
    box.fill.fore_color.rgb = PRIMARY
    box.line.fill.background()

    t = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(8.4), Inches(1.5))
    tf = t.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE

    if subtitle:
        s = slide.shapes.add_textbox(Inches(0.8), Inches(3.8), Inches(8.4), Inches(1))
        sp = s.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(18)
        sp.font.color.rgb = GRAY


def add_section_slide(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PRIMARY)
    t = slide.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(8.4), Inches(1.2))
    p = t.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER


def add_content_slide(prs, title, bullets, subtitle=''):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.1))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK
    header.line.fill.background()

    ht = slide.shapes.add_textbox(Inches(0.6), Inches(0.25), Inches(8.8), Inches(0.7))
    hp = ht.text_frame.paragraphs[0]
    hp.text = title
    hp.font.size = Pt(24)
    hp.font.bold = True
    hp.font.color.rgb = WHITE

    body = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(8.6), Inches(5.5))
    tf = body.text_frame
    tf.word_wrap = True

    if subtitle:
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(14)
        p.font.color.rgb = GRAY
        p.font.italic = True
        p.space_after = Pt(12)

    start_idx = 1 if subtitle else 0
    for i, bullet in enumerate(bullets):
        idx = start_idx + i
        if idx == 0 and not subtitle:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(16)
        p.font.color.rgb = DARK
        p.level = 0
        p.space_after = Pt(8)


def add_flow_box(slide, left, top, width, height, text, fill_color=PRIMARY, text_color=WHITE, font_size=11):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = fill_color
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = text_color
    return shape


def add_arrow_text(slide, left, top, text=''):
    if text:
        t = slide.shapes.add_textbox(left, top, Inches(0.4), Inches(0.3))
        p = t.text_frame.paragraphs[0]
        p.text = '▼'
        p.font.size = Pt(14)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER


def add_flowchart_slide(prs, title, steps, decision_steps=None):
    """Vertical flowchart with optional decision branches."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.0))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK
    header.line.fill.background()
    ht = slide.shapes.add_textbox(Inches(0.6), Inches(0.22), Inches(8.8), Inches(0.6))
    ht.text_frame.paragraphs[0].text = title
    ht.text_frame.paragraphs[0].font.size = Pt(22)
    ht.text_frame.paragraphs[0].font.bold = True
    ht.text_frame.paragraphs[0].font.color.rgb = WHITE

    decision_steps = decision_steps or set()
    y = Inches(1.25)
    box_h = Inches(0.55)
    box_w = Inches(3.2)
    x = Inches(3.4)

    for i, step in enumerate(steps):
        is_decision = i in decision_steps
        color = ORANGE if is_decision else PRIMARY
        shape_type = MSO_SHAPE.DIAMOND if is_decision else MSO_SHAPE.ROUNDED_RECTANGLE
        h = Inches(0.7) if is_decision else box_h

        shape = slide.shapes.add_shape(shape_type, x, y, box_w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = color
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = step
        p.font.size = Pt(10 if is_decision else 11)
        p.font.bold = True
        p.font.color.rgb = WHITE

        y += h + Inches(0.05)
        if i < len(steps) - 1:
            add_arrow_text(slide, x + box_w / 2 - Inches(0.2), y - Inches(0.05))
            y += Inches(0.22)


def add_er_slide(prs, title, entities):
    """Entity-relationship table slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.0))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK
    header.line.fill.background()
    ht = slide.shapes.add_textbox(Inches(0.6), Inches(0.22), Inches(8.8), Inches(0.6))
    ht.text_frame.paragraphs[0].text = title
    ht.text_frame.paragraphs[0].font.size = Pt(22)
    ht.text_frame.paragraphs[0].font.bold = True
    ht.text_frame.paragraphs[0].font.color.rgb = WHITE

    rows = len(entities) + 1
    cols = 4
    table_shape = slide.shapes.add_table(rows, cols, Inches(0.4), Inches(1.2), Inches(9.2), Inches(0.4 * rows))
    table = table_shape.table

    headers = ['Entity', 'Primary Key', 'Key Relationships', 'Scoped to Company']
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(10)
            p.font.color.rgb = WHITE

    for r, (entity, pk, rels, scoped) in enumerate(entities, 1):
        for c, val in enumerate([entity, pk, rels, scoped]):
            cell = table.cell(r, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9)
                p.font.color.rgb = DARK
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BG


def add_db_hub_slide(prs):
    """Visual hub diagram for Company-centric design."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.0))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK
    header.line.fill.background()
    ht = slide.shapes.add_textbox(Inches(0.6), Inches(0.22), Inches(8.8), Inches(0.6))
    ht.text_frame.paragraphs[0].text = 'Database Hub — Company-Centric Multi-Tenant Design'
    ht.text_frame.paragraphs[0].font.size = Pt(22)
    ht.text_frame.paragraphs[0].font.bold = True
    ht.text_frame.paragraphs[0].font.color.rgb = WHITE

    hub = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.8), Inches(2.8), Inches(2.4), Inches(1.4))
    hub.fill.solid()
    hub.fill.fore_color.rgb = PRIMARY
    hub.line.color.rgb = PRIMARY
    tf = hub.text_frame
    tf.paragraphs[0].text = 'COMPANY\n(Root Tenant)'
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    satellites = [
        (Inches(0.5), Inches(1.5), 'Branch', GREEN),
        (Inches(7.5), Inches(1.5), 'Department', GREEN),
        (Inches(0.5), Inches(4.5), 'User / Auth', ORANGE),
        (Inches(7.5), Inches(4.5), 'Employee', ORANGE),
        (Inches(0.3), Inches(3.2), 'Payroll', RED),
        (Inches(8.0), Inches(3.2), 'Leave', RED),
        (Inches(1.5), Inches(5.8), 'Attendance', GRAY),
        (Inches(6.0), Inches(5.8), 'Recruitment', GRAY),
        (Inches(3.0), Inches(1.3), 'Performance', GRAY),
        (Inches(5.5), Inches(1.3), 'Disciplinary', GRAY),
    ]

    for left, top, label, color in satellites:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.8), Inches(0.6))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.color.rgb = color
        p = box.text_frame.paragraphs[0]
        p.text = label
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 1. Title
    add_title_slide(
        prs,
        'HRMS Pro',
        'Enterprise Human Resource Management System\nSystem Design • Flowcharts • Database Architecture\nAugust 2026',
    )

    # 2. Agenda
    add_content_slide(prs, 'Agenda', [
        'System Overview & Architecture',
        'Authentication & Security Flow',
        'Business Process Flowcharts',
        '  → Recruitment, Leave, Payroll, Attendance',
        'Database Design & Entity Relationships',
        'Multi-Tenant Company Hub Model',
        'Table Reference (58+ entities)',
        'Technology Stack & Deployment',
    ])

    # 3. System Overview
    add_content_slide(
        prs,
        'System Overview',
        [
            'All-in-one HR & Payroll platform for SMB to Enterprise',
            '17 modular Django apps with API-first architecture',
            'Multi-company, multi-branch, multi-country support',
            'Consolidates: HR, Payroll, Leave, Recruitment, Attendance, Performance',
            'AI-powered: Resume screening, anomaly detection, attrition prediction, chatbot',
            'Enterprise security: RBAC, MFA, JWT, LDAP, OAuth, Audit logs',
        ],
        subtitle='Eliminates fragmented HR systems across departments',
    )

    # 4. Architecture
    add_section_slide(prs, 'Architecture')

    add_content_slide(
        prs,
        'Technology Stack',
        [
            'Backend: Django 5 + Django REST Framework',
            'Database: PostgreSQL 16 (SQLite for dev)',
            'Cache/Queue: Redis + Celery + Celery Beat',
            'Auth: JWT, django-allauth, django-otp (MFA)',
            'Frontend: Bootstrap 5, Chart.js, HTMX-ready',
            'Server: Gunicorn + Nginx reverse proxy',
            'Container: Docker Compose (web, db, redis, celery, nginx)',
            'CI/CD: GitHub Actions with pytest + ruff',
            'API Docs: OpenAPI / Swagger at /api/docs/',
        ],
    )

    add_content_slide(
        prs,
        'Application Architecture Flow',
        [
            'Client (Web / Mobile / API) → Nginx Gateway',
            'Nginx → Gunicorn (Django WSGI)',
            'Middleware Chain:',
            '  Security → CORS → Auth → OTP → Audit Log → Session Activity',
            'Route → ViewSet (DRF) → Service Layer → PostgreSQL',
            'Async Tasks → Celery Workers → Redis Broker',
            'Response: JSON (API) or HTML (Dashboard)',
        ],
    )

    # 5. Auth Flow
    add_section_slide(prs, 'Authentication Flow')

    add_flowchart_slide(
        prs,
        'Authentication Flow',
        [
            'User Login Request',
            'Select Auth Method',
            'Validate Credentials',
            'Account Locked?',
            'Password Valid?',
            'MFA Enabled?',
            'Verify TOTP Token',
            'Create Session + JWT',
            'Load RBAC Permissions',
            'Dashboard / API Access',
        ],
        decision_steps={3, 4, 5, 6},
    )

    add_content_slide(
        prs,
        'User Roles (RBAC)',
        [
            'Super Admin — Full system access',
            'HR Administrator — HR module management',
            'Payroll Officer — Payroll processing & approval',
            'Recruiter — Job postings, ATS, interviews',
            'Manager / Supervisor — Team approvals & oversight',
            'Employee — Self-service portal (ESS)',
            'Finance Officer — Financial reports & payroll approval',
            'Department Head — Department-scoped data',
            'Casual Supervisor — Casual worker management',
            'Auditor — Read-only access to audit logs',
        ],
    )

    # 6. Business Flows
    add_section_slide(prs, 'Business Process Flowcharts')

    add_flowchart_slide(
        prs,
        'Recruitment & Onboarding Flow',
        [
            'HR Creates Job Posting',
            'Publish to Career Portal',
            'Applicant Submits Resume',
            'AI Resume Screening & Ranking',
            'Qualified Candidate?',
            'Schedule & Conduct Interview',
            'Generate Offer Letter',
            'Offer Accepted?',
            'Create Employee Record',
            'Digital Onboarding Checklist',
            'Employee Active',
        ],
        decision_steps={4, 7},
    )

    add_flowchart_slide(
        prs,
        'Leave Management Flow',
        [
            'Employee Applies for Leave',
            'Check Leave Balance',
            'Sufficient Balance?',
            'Create LeaveRequest (Pending)',
            'Supervisor Approval (L1)',
            'Approved?',
            'Manager/HR Approval (L2)',
            'Final Approval → Deduct Balance',
            'Update Attendance Calendar',
        ],
        decision_steps={2, 5},
    )

    add_flowchart_slide(
        prs,
        'Payroll Processing Flow',
        [
            'Create Payroll Run (Draft)',
            'Celery: Process All Employees',
            'Calculate Gross Pay',
            'Apply Deductions & Tax',
            'Generate Payslips',
            'AI Anomaly Detection',
            'Payroll Officer Review',
            'Finance Approval',
            'Status: Paid → PDF Payslips',
        ],
        decision_steps={6},
    )

    add_flowchart_slide(
        prs,
        'Time & Attendance Flow',
        [
            'Employee Check-In',
            'Method: Manual / Biometric / QR / GPS / Face',
            'Create AttendanceRecord',
            'Compare with Shift & Roster',
            'Late Arrival?',
            'AI Anomaly Check',
            'Check-Out → Calculate Hours & Overtime',
            'Sync to Payroll Module',
        ],
        decision_steps={4},
    )

    add_flowchart_slide(
        prs,
        'Disciplinary Process Flow',
        [
            'Incident Reported',
            'Create Incident Record',
            'Investigation',
            'Assess Severity',
            'Verbal / Written Warning',
            'Disciplinary Hearing',
            'Suspension or Termination',
            'Case Closed + Audit Trail',
        ],
        decision_steps={3},
    )

    # 7. Database Design
    add_section_slide(prs, 'Database Design')

    add_db_hub_slide(prs)

    add_er_slide(prs, 'Core & Organization Entities', [
        ('Company', 'Auto ID', 'Root tenant — all entities', 'N/A (root)'),
        ('Branch', 'Auto ID', 'Company → Branch (1:N)', 'Yes'),
        ('Department', 'Auto ID', 'Company, Branch, Parent, Head (1:N)', 'Yes'),
        ('Designation', 'Auto ID', 'Company → Designation (1:N)', 'Yes'),
        ('Holiday', 'Auto ID', 'Company, Branch (1:N)', 'Yes'),
        ('SystemSetting', 'Auto ID', 'Company key-value config', 'Optional'),
    ])

    add_er_slide(prs, 'Authentication & Security Entities', [
        ('User', 'UUID', 'Company, Branch, Groups (1:N)', 'Yes'),
        ('PermissionGroup', 'Auto ID', 'Company, JSON permissions', 'Optional'),
        ('UserPermissionGroup', 'Auto ID', 'User ↔ Group (M:N)', '—'),
        ('APIToken', 'Auto ID', 'User → Token (1:N)', '—'),
        ('AuditLog', 'Auto ID', 'User, Company, action, changes', 'Optional'),
        ('UserSession', 'Auto ID', 'User session tracking (1:N)', '—'),
    ])

    add_er_slide(prs, 'Employee Management Entities', [
        ('Employee', 'UUID', 'Company, User(1:1), Manager(self)', 'Yes'),
        ('EmployeeContract', 'Auto ID', 'Employee → Contract (1:N)', 'Yes'),
        ('EmployeeDocument', 'Auto ID', 'Employee → Document (1:N)', 'Yes'),
        ('EmployeeHistory', 'Auto ID', 'Employee change audit (1:N)', '—'),
    ])

    add_er_slide(prs, 'Recruitment Entities', [
        ('JobPosting', 'UUID', 'Company, Dept, Designation (1:N)', 'Yes'),
        ('Applicant', 'UUID', 'Job → Applicant (1:N), AI score', 'Yes'),
        ('Interview', 'Auto ID', 'Applicant, Interviewer (1:N)', 'Yes'),
        ('OfferLetter', 'Auto ID', 'Applicant (1:1)', 'Yes'),
        ('OnboardingChecklist', 'Auto ID', 'Employee → Tasks (1:N)', 'Yes'),
    ])

    add_er_slide(prs, 'Payroll Entities', [
        ('SalaryStructure', 'Auto ID', 'Company templates (1:N)', 'Yes'),
        ('Allowance / Deduction', 'Auto ID', 'Company definitions (1:N)', 'Yes'),
        ('EmployeeSalary', 'Auto ID', 'Employee (1:1), M:N allowances', 'Yes'),
        ('PayrollRun', 'UUID', 'Company → Run (1:N), approved_by', 'Yes'),
        ('Payslip', 'Auto ID', 'PayrollRun, Employee (1:N)', 'Yes'),
        ('Loan', 'Auto ID', 'Employee → Loan (1:N)', 'Yes'),
    ])

    add_er_slide(prs, 'Leave & Attendance Entities', [
        ('LeaveType', 'Auto ID', 'Company leave policies (1:N)', 'Yes'),
        ('LeaveBalance', 'Auto ID', 'Employee + Type + Year (1:N)', 'Yes'),
        ('LeaveRequest', 'UUID', 'Employee, Type, Approvals (1:N)', 'Yes'),
        ('LeaveApproval', 'Auto ID', 'Request → Multi-level (1:N)', 'Yes'),
        ('Shift', 'Auto ID', 'Company shift definitions (1:N)', 'Yes'),
        ('AttendanceRecord', 'UUID', 'Employee + Date unique (1:N)', 'Yes'),
        ('Roster', 'Auto ID', 'Employee + Shift + Date (1:N)', 'Yes'),
        ('BiometricDevice', 'Auto ID', 'Branch device registry (1:N)', 'Yes'),
    ])

    add_er_slide(prs, 'Performance, Relations & Other Entities', [
        ('PerformanceCycle', 'UUID', 'Company review cycles (1:N)', 'Yes'),
        ('Goal / KPI', 'Auto ID', 'Employee goals, dept KPIs', 'Yes'),
        ('PerformanceReview', 'UUID', 'Employee + Cycle (1:N)', 'Yes'),
        ('Feedback360', 'Auto ID', 'Review → Feedback (1:N)', 'Yes'),
        ('Grievance / Recognition', 'UUID', 'Employee relations (1:N)', 'Yes'),
        ('Incident / Warning', 'UUID', 'Disciplinary chain (1:N)', 'Yes'),
        ('CasualWorker', 'UUID', 'Casual mgmt + attendance (1:N)', 'Yes'),
        ('Survey / Question / Response', 'UUID', 'Feedback surveys (1:N)', 'Yes'),
        ('AIAnalysisJob', 'UUID', 'AI job queue & results', 'Yes'),
        ('IntegrationProvider', 'Auto ID', 'ERP, SMS, payment gateways', 'Yes'),
        ('Notification', 'UUID', 'User notifications (1:N)', '—'),
    ])

    add_content_slide(
        prs,
        'Key Relationship Rules',
        [
            'Company is the root tenant — virtually all business data scoped via company_id FK',
            'Employee is the central HR entity — links to User (1:1), Manager (self-ref), Org structure',
            'User ↔ Employee: optional 1:1 link enabling Employee Self-Service (ESS)',
            'LeaveRequest → LeaveApproval: multi-level approval chain (level 1, 2, 3...)',
            'PayrollRun → Payslip: one run generates N payslips; AI flags anomalies per payslip',
            'JobPosting → Applicant → OfferLetter → Employee: recruitment-to-hire pipeline',
            'Incident → Warning → Hearing → Suspension: disciplinary escalation chain',
            'EmployeeSalary M:N Allowances & Deductions via junction tables',
            'All tables include created_at / updated_at timestamps (TimeStampedModel)',
        ],
        subtitle='58+ custom tables + Django system tables',
    )

    # 8. Closing
    add_content_slide(
        prs,
        'Deployment Architecture',
        [
            'Docker Compose services: web, db, redis, celery, celery-beat, nginx',
            'Production: AWS / Azure with managed PostgreSQL & Redis',
            'Environment config via .env (django-environ)',
            'Static files: WhiteNoise + Nginx',
            'Background jobs: payroll processing, email, leave accrual, AI screening',
            'Backup & Restore: PostgreSQL pg_dump + media files',
            'Monitoring: Sentry integration (production)',
        ],
    )

    add_title_slide(
        prs,
        'Thank You',
        'HRMS Pro v1.0.0\nDocumentation: docs/HRMS-System-Design-Documentation.md\nAPI Docs: /api/docs/',
    )

    prs.save(str(OUTPUT))
    print(f'Presentation saved to: {OUTPUT}')


if __name__ == '__main__':
    build_presentation()
